# -*- coding: utf-8 -*-
"""Сборка презентации Amber в PPTX и экспорт в PDF (PowerPoint COM или LibreOffice)."""
from __future__ import annotations

import os
import subprocess
import sys

from pptx import Presentation
from pptx.util import Inches, Pt


DIR = os.path.dirname(os.path.abspath(__file__))
OUT_PPTX = os.path.join(DIR, "amber-ai-compliance-copilot.pptx")
OUT_PDF = os.path.join(DIR, "amber-ai-compliance-copilot.pdf")
OUT_HTML = os.path.join(DIR, "amber-ai-compliance-copilot.html")


def _set_slide_size_wide(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)


def _font_rt(run, size_pt: float, bold: bool = False) -> None:
    run.font.name = "Calibri"
    run.font.size = Pt(size_pt)
    run.font.bold = bold


def add_title_slide(prs: Presentation, title: str, subtitle_lines: list[str]) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(subtitle_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        if p.runs:
            _font_rt(p.runs[0], 16, False)


def add_title_content(prs: Presentation, title: str, body_lines: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(body_lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        if p.runs:
            _font_rt(p.runs[0], 15, False)


def add_title_mono_block(prs: Presentation, title: str, mono_text: str, footer: str | None = None) -> None:
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    left, top = Inches(0.6), Inches(0.45)
    tw = Inches(12.1)
    title_box = slide.shapes.add_textbox(left, top, tw, Inches(0.85))
    title_box.text_frame.text = title
    for r in title_box.text_frame.paragraphs[0].runs:
        _font_rt(r, 26, True)

    body_top = Inches(1.35)
    body_h = Inches(5.35) if footer else Inches(5.9)
    box = slide.shapes.add_textbox(left, body_top, tw, body_h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in mono_text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        if p.runs:
            r = p.runs[0]
            r.font.name = "Consolas"
            r.font.size = Pt(9)
            r.font.bold = False

    if footer:
        fb = slide.shapes.add_textbox(left, Inches(6.75), tw, Inches(0.65))
        fb.text_frame.text = footer
        for r in fb.text_frame.paragraphs[0].runs:
            _font_rt(r, 11, False)


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    footer: str | None = None,
) -> None:
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.38), Inches(12.1), Inches(0.8))
    tbox.text_frame.text = title
    for r in tbox.text_frame.paragraphs[0].runs:
        _font_rt(r, 26, True)

    nrows = 1 + len(rows)
    ncols = len(headers)
    row_h = min(0.52, 5.8 / max(nrows, 4))
    tbl_top = Inches(1.25)
    tbl_h = Inches(row_h * nrows + 0.25)
    tbl = slide.shapes.add_table(nrows, ncols, Inches(0.55), tbl_top, Inches(12.2), tbl_h).table
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                _font_rt(run, 11, True)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    _font_rt(run, 10, False)

    if footer:
        ft = slide.shapes.add_textbox(Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.75))
        ft.text_frame.word_wrap = True
        ft.text_frame.text = footer
        for r in ft.text_frame.paragraphs[0].runs:
            _font_rt(r, 12, False)


def export_pdf_powerpoint(pptx_path: str, pdf_path: str) -> bool:
    try:
        import win32com.client  # type: ignore
    except ImportError:
        return False
    pptx_path = os.path.abspath(pptx_path)
    pdf_path = os.path.abspath(pdf_path)
    app = None
    pres = None
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        app.Visible = 0
        pres = app.Presentations.Open(pptx_path, WithWindow=False)
        pres.SaveAs(pdf_path, 32)
        pres.Close()
        pres = None
        app.Quit()
        app = None
        return os.path.isfile(pdf_path)
    except Exception:
        try:
            if pres:
                pres.Close()
        except Exception:
            pass
        try:
            if app:
                app.Quit()
        except Exception:
            pass
        return False


def _find_chromium_exe() -> str | None:
    candidates = [
        os.path.join(os.environ.get("ProgramFiles", ""), "Google", "Chrome", "Application", "chrome.exe"),
        os.path.join(os.environ.get("ProgramFiles(x86)", ""), "Google", "Chrome", "Application", "chrome.exe"),
        os.path.join(os.environ.get("LocalAppData", ""), "Google", "Chrome", "Application", "chrome.exe"),
        os.path.join(os.environ.get("ProgramFiles(x86)", ""), "Microsoft", "Edge", "Application", "msedge.exe"),
        os.path.join(os.environ.get("ProgramFiles", ""), "Microsoft", "Edge", "Application", "msedge.exe"),
    ]
    for c in candidates:
        if c and os.path.isfile(c):
            return c
    return None


def export_pdf_from_html_chromium(html_path: str, pdf_path: str) -> bool:
    browser = _find_chromium_exe()
    if not browser or not os.path.isfile(html_path):
        return False
    html_abs = os.path.abspath(html_path)
    pdf_abs = os.path.abspath(pdf_path)
    # file:/// URL with forward slashes
    url = "file:///" + html_abs.replace("\\", "/")
    try:
        subprocess.run(
            [
                browser,
                "--headless=new",
                "--disable-gpu",
                "--no-pdf-header-footer",
                f"--print-to-pdf={pdf_abs}",
                url,
            ],
            check=True,
            capture_output=True,
            text=True,
            timeout=120,
        )
        return os.path.isfile(pdf_abs) and os.path.getsize(pdf_abs) > 1000
    except Exception:
        return False


def export_pdf_libreoffice(pptx_path: str, pdf_path: str) -> bool:
    candidates = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    soffice = next((c for c in candidates if os.path.isfile(c)), None)
    if not soffice:
        return False
    outdir = os.path.dirname(os.path.abspath(pdf_path))
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", outdir, os.path.abspath(pptx_path)],
            check=True,
            capture_output=True,
            text=True,
            timeout=120,
        )
        produced = os.path.join(outdir, base + ".pdf")
        if os.path.isfile(produced) and os.path.abspath(produced) != os.path.abspath(pdf_path):
            if os.path.isfile(pdf_path):
                os.remove(pdf_path)
            os.replace(produced, pdf_path)
        return os.path.isfile(pdf_path)
    except Exception:
        return False


def build() -> Presentation:
    prs = Presentation()
    _set_slide_size_wide(prs)

    add_title_slide(
        prs,
        "Amber — AI Compliance Copilot",
        [
            "Интеллектуальный помощник для compliance-офицера белорусского банка",
            "",
            "Минск, 2026",
            "",
            "Контекст: банковская система Беларуси",
            "• С 1 июля 2026 — единые стандарты антифрод-систем: автоматический анализ и блокировка подозрительных операций.",
            "• Криптобанки (Указ №19): фиат и крипта на одном балансе; классический AML не видит связку.",
            "• 23 крипторезидента ПВТ — легальные фиатные платежи через банковскую систему.",
            "• Мало AML-специалистов; объём алертов растёт.",
        ],
    )

    add_title_content(
        prs,
        "Проблема: текущие AML-системы",
        [
            "«Чёрный ящик» — флаг Suspicious без объяснения; офицер часами разбирается.",
            "95% ложных срабатываний — жёсткие пороги дают лавину пустых алертов.",
            "Слепая зона фиат ↔ крипта — сквозной картины нет (фиат в банке + крипта на санкционный кошелёк).",
            "SAR вручную — 1,5–2 часа на документ, сотни часов в месяц.",
        ],
    )

    diagram = """  Ваша AML-система                Amber API
  ┌──────────────────┐           ┌──────────────────┐
  │  Алерт           │──────────▶│  POST /analyze   │
  │  Экспорт (CSV)   │           └────────┬─────────┘
  └──────────────────┘                    │
                                          ▼
                                ┌──────────────────┐
                                │  XAI Engine      │
                                │  Profiler / AnomalyDet │
                                │  Analyst / Reporter    │
                                └────────┬─────────┘
                                         ▼
                                Ответ за 3–5 с: объяснение, SAR, аномалии"""
    add_title_mono_block(
        prs,
        "Решение: Amber — слой поверх AML",
        diagram,
        "Подключение — 1 рабочий день. Без замены AML. Без доступа к вашим БД.",
    )

    add_title_content(
        prs,
        "Что Amber даёт банку",
        [
            "1. Объяснимый ИИ — паттерн, причина риска, норма, что делать дальше.",
            "2. Автоматический проект SAR для ДФР (русский, формат РБ, Декрет №8, Указ №19, Нацбанк).",
            "3. Аномалии вне правил — профиль «нормы» клиента и отклонения.",
        ],
    )

    add_table_slide(
        prs,
        "Три режима работы",
        ["Режим", "Для чего", "Что находит"],
        [
            ["fiat", "Классический банковский AML", "Структурирование, смурфинг, круговые переводы, профиль, аномалии"],
            ["crypto", "Крипто-операции", "Санкции, миксеры, PEP"],
            ["cross", "Связка фиат ↔ крипта", "Временные и суммовые корреляции фиат ↔ крипта"],
        ],
        footer="Юрисдикции: РБ (Декрет №8, Указ №19, Нацбанк), РФ (115-ФЗ), ЕС (5AMLD + MiCA).",
    )

    add_title_content(
        prs,
        "Детектор аномалий",
        [
            "Обычный AML: «сумма > порога → флаг».",
            "Amber: профиль за 90 дней — сравнение с нормой и объяснение.",
            "Шаги: Profiler строит норму (чек, частота, контрагенты, география, время); AnomalyDetector сравнивает; anomaly_score 0–100 + текст.",
            "Пример: score 65/100 — сумма +189% к историческому максимуму; 7 операций при норме до 4; новый контрагент ООО «Технопром»; 23:15.",
            "Гипотеза: дробление новому контрагенту вне часов — проверить ЕГРЮЛ/ФНС.",
            "Без ML — статистика, прозрачно и воспроизводимо.",
        ],
    )

    sar_block = """Сообщение о подозрительной операции № AMBER-2026-0506-001

1. Резюме: доход 1 000 BYN/мес; за 3 дня внесено 17 650 BYN частями 5 800–5 950 BYN — признаки structuring.

2. Операции:
| 29.04.2026 | 5 800 | Внесение |
| 30.04.2026 | 5 900 | Внесение |
| 02.05.2026 | 5 950 | Внесение |

3. Паттерны: Structuring; Income mismatch (оборот в 5.8× выше дохода).

4. Нормы: Декрет №8 ст. 7; Правила Нацбанка РБ о внутреннем контроле.

5. Рекомендации: запрос происхождения средств; при отказе — эскалация в ДФР.

⚠ Сгенерировано Amber AI — проверка compliance-офицером."""
    add_title_mono_block(prs, "Пример SAR от Amber", sar_block, None)

    arch = """┌──────────────────────────────────────────────────────┐
│                FASTAPI (main.py)                     │
│  POST /analyze  ·  GET /health  ·  CORS              │
└────────────────────┬─────────────────────────────────┘
                     │
┌────────────────────▼─────────────────────────────────┐
│               XAI ENGINE                           │
│  Profiler · AnomalyDetector · Router · Analyst · Reporter │
│  Ответы LLM — Pydantic-валидация                   │
└────────────────────┬─────────────────────────────────┘
                     │
┌────────────────────▼─────────────────────────────────┐
│  LLM: GPT-4o → Claude (резерв) · Emergency Mode     │
│  Zero retention                                    │
└──────────────────────────────────────────────────────┘

Принципы: stateless; без хранения данных; логи без ПДн; Docker в контуре банка."""
    add_title_mono_block(prs, "Техническая архитектура", arch, None)

    add_table_slide(
        prs,
        "Безопасность и соответствие",
        ["Аспект", "Как обеспечено"],
        [
            ["Хранение данных", "Только RAM на время запроса"],
            ["LLM API", "Zero retention (OpenAI, Anthropic)"],
            ["Логирование", "latency, статус, токены — без ПДн"],
            ["Инфраструктура", "ЦОД банка или beCloud"],
            ["Отказоустойчивость", "Два LLM + Emergency Mode"],
        ],
    )

    add_title_content(
        prs,
        "Интеграция — 1 день",
        [
            "Шаг 1: экспорт алерта (CSV/JSON). Шаг 2: загрузка в Amber (API или форма). Шаг 3: результат за 3–5 с.",
            "Результат: объяснение риска, проект SAR, anomaly score, рекомендации.",
            "Не нужно: замена AML, доступ к БД, долгое обучение, полугодовой цикл внедрения.",
        ],
    )

    add_title_content(
        prs,
        "Почему Amber — сейчас",
        [
            "1. Регулятор: с 01.07.2026 единая антифрод-система с ИИ — Amber закрывает требование автоматического анализа.",
            "2. Крипто-трансформация (Указ №19): без сквозного анализа — повышенный риск.",
            "3. Кадры: офицер с Amber обрабатывает ~в 3 раза больше алертов.",
            "4. Локально под РБ/ПВТ; глобальные вендоры не заточены под двойной надзор.",
        ],
    )

    add_title_content(
        prs,
        "Пилотное предложение",
        [
            "Формат: бесплатный пилот на 50 анализов, срок 2 недели.",
            "От банка: обезличенные данные по 10–20 реальным алертам.",
            "Вы получаете: полный разбор + SAR; сравнение времени ручного разбора vs Amber; отчёт по паттернам; рекомендации по интеграции.",
            "Итог: решение о внедрении на ваших цифрах.",
        ],
    )

    add_title_content(
        prs,
        "Приложение: технические детали",
        [
            "API: REST, JSON · HTTPS · API-ключ · CSV/JSON · ответ 3–5 с · Docker.",
            "Стек: Python 3.12, FastAPI, GPT-4o / Claude 3.5 Sonnet, LangChain, pandas, Pydantic.",
            "Инфраструктура банка: HTTPS к внешнему API или Docker в контуре.",
        ],
    )

    return prs


def main() -> int:
    prs = build()
    prs.save(OUT_PPTX)
    print("Записано:", OUT_PPTX)

    ok = export_pdf_powerpoint(OUT_PPTX, OUT_PDF)
    if ok:
        print("PDF (PowerPoint):", OUT_PDF)
        return 0

    ok = export_pdf_libreoffice(OUT_PPTX, OUT_PDF)
    if ok:
        print("PDF (LibreOffice):", OUT_PDF)
        return 0

    if os.path.isfile(OUT_HTML):
        ok = export_pdf_from_html_chromium(OUT_HTML, OUT_PDF)
        if ok:
            print("PDF (Chrome/Edge, из HTML-слайдов):", OUT_PDF)
            return 0

    print(
        "PDF не создан автоматически. Установите LibreOffice, либо Chrome/Edge, либо откройте "
        "amber-ai-compliance-copilot.pptx в PowerPoint → Сохранить как PDF.",
        file=sys.stderr,
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
